#!/usr/bin/env python3
"""
ERP System PowerPoint Presentation Generator
Creates a comprehensive PowerPoint presentation showcasing all ERP system modules and functions.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_erp_presentation():
    """Create a comprehensive ERP system PowerPoint presentation."""
    
    # Create presentation object
    prs = Presentation()
    
    # Define colors
    primary_blue = RGBColor(25, 118, 210)  # #1976d2
    secondary_blue = RGBColor(66, 165, 245)  # #42a5f5
    text_gray = RGBColor(51, 51, 51)  # #333333
    light_gray = RGBColor(248, 249, 250)  # #f8f9fa
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "ERP System"
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.color.rgb = primary_blue
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Complete Business Management Solution\nWeb & Mobile Application\n\nComprehensive Enterprise Resource Planning Platform"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[1].font.size = Pt(18)
    subtitle.text_frame.paragraphs[2].font.size = Pt(16)
    
    # Slide 2: System Overview
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "System Overview"
    title.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content = slide.placeholders[1]
    content.text = """A comprehensive ERP solution designed to streamline business operations across all departments with modern web and mobile interfaces.

Key Statistics:
• 12+ Core Modules
• 50+ Features  
• 100% Mobile Responsive
• 24/7 Availability

Technology Stack:
• React.js - Modern frontend framework
• Django - Robust backend API
• PostgreSQL - Reliable database
• Material-UI - Professional design
• Capacitor - Cross-platform mobile
• JWT Auth - Secure authentication"""
    
    # Slide 3: Core Features
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Core Features"
    title.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content = slide.placeholders[1]
    content.text = """🔐 Secure Authentication
JWT-based authentication with role-based access control and password reset functionality

📱 Mobile Ready
Fully responsive design with native mobile app support for iOS and Android

📊 Real-time Analytics
Live dashboards with interactive charts and comprehensive reporting

🔄 Workflow Management
Automated approval workflows with email and SMS notifications"""
    
    # Module slides data
    modules = [
        {
            "title": "Main Dashboard",
            "icon": "📊",
            "features": [
                "Executive summary cards",
                "Real-time KPI monitoring", 
                "Interactive charts and graphs",
                "Quick action buttons",
                "Recent activity feed",
                "Module navigation",
                "Role-based dashboard views",
                "Customizable widgets"
            ]
        },
        {
            "title": "HR Management Module",
            "icon": "👥",
            "features": [
                "Employee management",
                "Leave request system",
                "Payroll processing",
                "Recruitment & onboarding",
                "Training management",
                "HR calendar & events",
                "Department management",
                "Performance tracking"
            ]
        },
        {
            "title": "Inventory Management",
            "icon": "📦",
            "features": [
                "Product catalog management",
                "Stock level monitoring",
                "Warehouse management",
                "Stock transfers",
                "Inventory valuation",
                "Reorder point alerts",
                "Barcode scanning",
                "Stock movement tracking"
            ]
        },
        {
            "title": "Sales Management",
            "icon": "💰",
            "features": [
                "Sales order management",
                "Customer relationship management",
                "Quote generation",
                "Sales analytics",
                "Commission tracking",
                "Territory management",
                "Sales pipeline",
                "Revenue forecasting"
            ]
        },
        {
            "title": "Point of Sale (POS)",
            "icon": "🛒",
            "features": [
                "Touch-friendly interface",
                "Product scanning",
                "Multiple payment methods",
                "Receipt printing",
                "Daily sales reports",
                "Cash drawer management",
                "Customer lookup",
                "Discount management"
            ]
        },
        {
            "title": "Accounting & Finance",
            "icon": "💳",
            "features": [
                "General ledger",
                "Accounts payable/receivable",
                "Financial reporting",
                "Budget management",
                "Tax calculations",
                "Bank reconciliation",
                "Cash flow analysis",
                "Audit trails"
            ]
        },
        {
            "title": "Procurement Management",
            "icon": "📋",
            "features": [
                "Purchase requisitions",
                "Purchase orders",
                "Vendor management",
                "Approval workflows",
                "Bid analysis",
                "Contract management",
                "Supplier evaluation",
                "Cost analysis"
            ]
        },
        {
            "title": "Manufacturing Module",
            "icon": "🏭",
            "features": [
                "Production planning",
                "Work order management",
                "Bill of materials (BOM)",
                "Quality control",
                "Resource scheduling",
                "Production tracking",
                "Capacity planning",
                "Cost accounting"
            ]
        },
        {
            "title": "Customer Management",
            "icon": "👥",
            "features": [
                "Customer database",
                "Contact management",
                "Customer segmentation",
                "Communication history",
                "Credit management",
                "Customer analytics",
                "Loyalty programs",
                "Support tickets"
            ]
        },
        {
            "title": "Reporting & Analytics",
            "icon": "📈",
            "features": [
                "Custom report builder",
                "Interactive dashboards",
                "Data visualization",
                "Scheduled reports",
                "Export capabilities",
                "KPI monitoring",
                "Trend analysis",
                "PowerBI integration"
            ]
        },
        {
            "title": "Mobile Application",
            "icon": "📱",
            "features": [
                "Native iOS & Android apps",
                "Offline synchronization",
                "Push notifications",
                "Touch-optimized interface",
                "Camera integration",
                "GPS location services",
                "Biometric authentication",
                "Real-time data sync"
            ]
        },
        {
            "title": "User Management & Security",
            "icon": "🔐",
            "features": [
                "Role-based access control",
                "User permissions",
                "Department assignments",
                "Password policies",
                "Session management",
                "Audit logging",
                "Two-factor authentication",
                "Data encryption"
            ]
        },
        {
            "title": "Workflow & Notifications",
            "icon": "🔄",
            "features": [
                "Approval workflows",
                "Email notifications",
                "SMS alerts",
                "Push notifications",
                "Task assignments",
                "Deadline tracking",
                "Status monitoring",
                "Escalation rules"
            ]
        }
    ]
    
    # Create module slides
    for module in modules:
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"{module['icon']} {module['title']}"
        title.text_frame.paragraphs[0].font.color.rgb = primary_blue
        
        content = slide.placeholders[1]
        features_text = "\n".join([f"• {feature}" for feature in module['features']])
        content.text = f"Key Features:\n\n{features_text}\n\n[Screenshot placeholder: {module['title']} interface showing key functionality and user experience]"
    
    # Technical Architecture Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Technical Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content = slide.placeholders[1]
    content.text = """⚛️ Frontend
React.js with Material-UI for modern, responsive user interfaces

🐍 Backend  
Django REST Framework for robust API development

🗄️ Database
PostgreSQL for reliable data storage and management

📱 Mobile
Capacitor for cross-platform mobile app development

🔒 Security
JWT authentication, role-based access, data encryption

☁️ Deployment
Cloud-ready architecture with scalable infrastructure"""
    
    # Benefits Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Key Benefits"
    title.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content = slide.placeholders[1]
    content.text = """⚡ Increased Efficiency
Streamline operations and reduce manual processes by up to 70%

💰 Cost Reduction
Lower operational costs through automation and optimization

📊 Better Insights
Real-time analytics for informed decision making

🔄 Scalability
Grows with your business needs and requirements

🛡️ Enhanced Security
Enterprise-grade security with role-based access control

📱 Mobile Accessibility
Work from anywhere with full mobile functionality"""
    
    # Implementation & Support Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Implementation & Support"
    title.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content = slide.placeholders[1]
    content.text = """Implementation Process:
• Requirements analysis
• System configuration
• Data migration
• User training
• Go-live support
• Post-implementation review
• Ongoing maintenance
• 24/7 technical support

Training & Support:
• Comprehensive user manuals
• Video tutorials
• Live training sessions
• Help desk support
• System updates
• Performance monitoring
• Backup & recovery
• Security updates"""
    
    # Final Slide - Contact & Demo
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Get Started Today"
    title.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content = slide.placeholders[1]
    content.text = """Ready to Transform Your Business?
Experience the power of our comprehensive ERP solution

Demo Credentials:
• Admin: admin@erp.com (password: admin123)
• Manager: manager@erp.com (password: manager123)  
• Employee: employee@erp.com (password: employee123)
• Sales: sales@erp.com (password: sales123)

Contact Information:
📧 Email: info@erpsystem.com
📞 Phone: +1 (555) 123-4567
🌐 Website: www.erpsystem.com

Schedule a personalized demo today and see how our ERP system can transform your business operations!"""
    
    return prs

def main():
    """Main function to create and save the presentation."""
    print("Creating ERP System PowerPoint Presentation...")
    
    try:
        # Create the presentation
        presentation = create_erp_presentation()
        
        # Save the presentation
        output_file = "ERP_System_Complete_Presentation.pptx"
        presentation.save(output_file)
        
        print(f"✅ Presentation created successfully: {output_file}")
        print(f"📊 Total slides: {len(presentation.slides)}")
        print("\n📋 Presentation includes:")
        print("• Title slide with system overview")
        print("• Core features and benefits")
        print("• Detailed module breakdowns (12+ modules)")
        print("• Technical architecture")
        print("• Implementation and support information")
        print("• Demo credentials and contact details")
        print("\n🎯 Ready for download and presentation!")
        
    except ImportError:
        print("❌ Error: python-pptx library not found.")
        print("📦 Please install it using: pip install python-pptx")
        print("\n📝 Alternative: Use the HTML file (ERP_System_Presentation.html)")
        print("   - Open in browser and print to PDF")
        print("   - Or convert using online HTML to PowerPoint tools")
        
    except Exception as e:
        print(f"❌ Error creating presentation: {str(e)}")
        print("\n📝 Alternative: Use the HTML file (ERP_System_Presentation.html)")

if __name__ == "__main__":
    main()
